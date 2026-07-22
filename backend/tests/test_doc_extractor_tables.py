"""
test_doc_extractor_tables.py — Unit tests for table extraction, Markdown formatting, document order, and merged cell handling.
"""

import os
import tempfile
import pytest
from services.doc_extractor import (
    _format_table_matrix_to_markdown,
    _extract_docx,
    _extract_pptx,
    extract_text_from_file,
)


def test_format_table_matrix_to_markdown_basic():
    """Verify matrix formatting generates structured Markdown with header delimiters and pipe escaping."""
    matrix = [
        ["Name", "Role", "Department | Team"],
        ["Alice", "Engineer", "Backend"],
        ["Bob", "Manager\nLead", "Product"],
    ]
    md = _format_table_matrix_to_markdown(matrix, title="Test Table")

    assert "[Test Table: 3 rows x 3 cols]" in md
    assert "| Name | Role | Department \\| Team |" in md
    assert "| --- | --- | --- |" in md
    assert "| Alice | Engineer | Backend |" in md
    assert "| Bob | Manager<br>Lead | Product |" in md


def test_format_table_matrix_to_markdown_padding():
    """Verify rows of uneven lengths are padded to the maximum column count."""
    matrix = [
        ["Col1", "Col2", "Col3"],
        ["Val1"],
        ["ValA", "ValB", "ValC"],
    ]
    md = _format_table_matrix_to_markdown(matrix, title="Padded Table")
    assert "| Val1 |  |  |" in md


def test_docx_table_and_document_flow_extraction():
    """Verify DOCX paragraphs and tables are extracted in natural document flow order."""
    import docx

    doc = docx.Document()
    doc.add_paragraph("First Paragraph before Table")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Data 1"
    table.cell(1, 1).text = "Data 2"

    doc.add_paragraph("Second Paragraph after Table")

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp_path = tmp.name
        doc.save(tmp_path)

    try:
        text = _extract_docx(tmp_path)

        pos_para1 = text.find("First Paragraph before Table")
        pos_table = text.find("[Table: 2 rows x 2 cols]")
        pos_para2 = text.find("Second Paragraph after Table")

        assert pos_para1 != -1
        assert pos_table != -1
        assert pos_para2 != -1

        # Verify document flow order: Para 1 < Table < Para 2
        assert pos_para1 < pos_table < pos_para2
        assert "| Header 1 | Header 2 |" in text
        assert "| Data 1 | Data 2 |" in text
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


def test_docx_merged_cells_extraction():
    """Verify horizontally merged cells in DOCX do not produce duplicated text across columns."""
    import docx

    doc = docx.Document()
    table = doc.add_table(rows=2, cols=3)

    # Merge top row across all 3 columns
    merged_header = table.cell(0, 0).merge(table.cell(0, 2))
    merged_header.text = "Spanning Header"

    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    table.cell(1, 2).text = "C"

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp_path = tmp.name
        doc.save(tmp_path)

    try:
        text = _extract_docx(tmp_path)
        assert "[Table: 2 rows x 3 cols]" in text
        assert "| Spanning Header |  |  |" in text
        assert "| A | B | C |" in text
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


def test_pptx_table_extraction():
    """Verify PPTX slide tables are extracted into Markdown tables with empty cells preserved."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Add a table
    table_shape = slide.shapes.add_table(rows=2, cols=3, left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(2))
    tbl = table_shape.table

    tbl.cell(0, 0).text = "Col A"
    tbl.cell(0, 1).text = "Col B"
    tbl.cell(0, 2).text = "Col C"

    tbl.cell(1, 0).text = "Val 1"
    tbl.cell(1, 1).text = ""  # Empty cell
    tbl.cell(1, 2).text = "Val 3"

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
        prs.save(tmp_path)

    try:
        text = _extract_pptx(tmp_path)
        assert "[Slide 1]" in text
        assert "[Table (Slide 1): 2 rows x 3 cols]" in text
        assert "| Col A | Col B | Col C |" in text
        assert "| Val 1 |  | Val 3 |" in text
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


def test_extract_text_from_file_supported_extensions():
    """Test extract_text_from_file dispatcher for docx and pptx files."""
    import docx

    doc = docx.Document()
    doc.add_paragraph("Sample Document Test")

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp_path = tmp.name
        doc.save(tmp_path)

    try:
        text = extract_text_from_file(tmp_path, "sample.docx")
        assert "Sample Document Test" in text
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
