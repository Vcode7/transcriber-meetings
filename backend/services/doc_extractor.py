"""
doc_extractor.py — Offline text extraction from uploaded agenda/context files.

Supported formats:
  - PDF   : PyMuPDF (fitz) text extraction + PyMuPDF embedded image extraction + Tesseract OCR
  - DOCX  : python-docx text extraction + inline image extraction + Tesseract OCR
  - PPTX  : python-pptx text extraction + picture shape extraction + Tesseract OCR
  - TXT   : plain read
  - MD    : plain read
  - PNG / JPG / JPEG / WEBP : pytesseract OCR
  - XLSX / XLS / CSV        : pandas spreadsheet sheet-by-sheet extraction + row indices tracking

Tesseract is always invoked with CREATE_NO_WINDOW so no console/cmd window
flickers in the packaged (.exe) application.  The pytesseract subprocess flags
are applied once at import time via _configure_tesseract_subprocess().
"""
from __future__ import annotations

import logging
import os
import re
import io
import subprocess
import sys
from typing import Optional

logger = logging.getLogger(__name__)

# ── Windows silent-subprocess constants ────────────────────────────────────────
# CREATE_NO_WINDOW (0x08000000) prevents a cmd/console window from appearing
# when tesseract.exe is spawned from a windowed / packaged Python process.
_CREATE_NO_WINDOW = 0x08000000

# ── Tesseract availability (cached after first check) ─────────────────────────
_TESSERACT_AVAILABLE: Optional[bool] = None
_TESSERACT_CONFIGURED: bool = False


def _get_bundled_tesseract_path() -> Optional[str]:
    """
    Return the path to a bundled tesseract.exe when running inside a
    PyInstaller-frozen executable, or None if running in development.

    Convention: the installer places Tesseract-OCR inside
      <app_root>/runtime/tesseract/tesseract.exe
    where <app_root> is the directory that contains backend.exe / launcher.exe.
    """
    if not getattr(sys, "frozen", False):
        return None  # development mode — rely on system PATH

    # sys.executable → …/Application/backend/backend.exe
    # app_root       → …/Application/
    try:
        exe_dir = os.path.dirname(sys.executable)
        app_root = os.path.dirname(exe_dir)
        candidate = os.path.join(app_root, "runtime", "tesseract", "tesseract.exe")
        if os.path.isfile(candidate):
            return candidate
    except Exception:
        pass
    return None


def _configure_tesseract_subprocess() -> None:
    """
    Patch pytesseract's internal subprocess call once so that:
      1. The correct tesseract.exe is used in the packaged app.
      2. No console/cmd window ever appears (CREATE_NO_WINDOW + SW_HIDE).

    This must be called before the first pytesseract.image_to_string() call.
    It is idempotent — repeated calls are a no-op.
    """
    global _TESSERACT_CONFIGURED
    if _TESSERACT_CONFIGURED:
        return

    try:
        import pytesseract

        # ── 1. Point to bundled binary if available ────────────────────────────
        bundled = _get_bundled_tesseract_path()
        if bundled:
            pytesseract.pytesseract.tesseract_cmd = bundled
            logger.info(f"[DocExtractor] Using bundled Tesseract: {bundled}")

        # ── 2. Patch subprocess so no window ever appears ─────────────────────
        # pytesseract.run_tesseract() uses subprocess.Popen internally.
        # We monkey-patch the module-level Popen reference so that every call
        # automatically carries CREATE_NO_WINDOW + STARTF_USESHOWWINDOW.
        if sys.platform == "win32":
            _orig_popen = subprocess.Popen

            class _SilentPopen(_orig_popen):  # type: ignore[misc]
                """Subprocess.Popen subclass that always hides the console window."""

                def __init__(self, *args, **kwargs):
                    # Build STARTUPINFO
                    si = kwargs.pop("startupinfo", None) or subprocess.STARTUPINFO()
                    si.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                    si.wShowWindow = subprocess.SW_HIDE
                    kwargs["startupinfo"] = si

                    # Merge creation flags
                    flags = kwargs.pop("creationflags", 0)
                    kwargs["creationflags"] = flags | _CREATE_NO_WINDOW

                    super().__init__(*args, **kwargs)

            # Patch only the subprocess reference inside the pytesseract module
            # so we do not affect the rest of the application.
            import pytesseract.pytesseract as _pt_core
            _pt_core.subprocess.Popen = _SilentPopen  # type: ignore[attr-defined]
            logger.debug("[DocExtractor] Patched pytesseract to use CREATE_NO_WINDOW")

        _TESSERACT_CONFIGURED = True

    except Exception as exc:
        logger.warning(f"[DocExtractor] Could not configure Tesseract subprocess: {exc}")


def _check_tesseract() -> bool:
    """Return True if Tesseract is installed and callable (result is cached)."""
    global _TESSERACT_AVAILABLE
    if _TESSERACT_AVAILABLE is not None:
        return _TESSERACT_AVAILABLE
    try:
        _configure_tesseract_subprocess()
        import pytesseract
        pytesseract.get_tesseract_version()
        _TESSERACT_AVAILABLE = True
        logger.info("[DocExtractor] Tesseract OCR is available")
    except Exception as exc:
        _TESSERACT_AVAILABLE = False
        logger.warning(f"[DocExtractor] Tesseract OCR not available: {exc}")
    return _TESSERACT_AVAILABLE


def _ocr_image_bytes(image_bytes: bytes, *, label: str = "<bytes>") -> str:
    """
    Run pytesseract OCR on in-memory image bytes.

    Retries at most once on transient failure.  If the retry also fails,
    the image is silently skipped and an empty string is returned.
    """
    if not _check_tesseract():
        return ""

    max_attempts = 2
    last_exc: Optional[Exception] = None

    for attempt in range(1, max_attempts + 1):
        try:
            from PIL import Image
            import pytesseract
            _configure_tesseract_subprocess()
            img = Image.open(io.BytesIO(image_bytes))
            return pytesseract.image_to_string(img).strip()
        except Exception as exc:
            last_exc = exc
            if attempt < max_attempts:
                logger.warning(
                    f"[DocExtractor] OCR attempt {attempt}/{max_attempts} failed for {label}: {exc} — retrying"
                )
            else:
                logger.warning(
                    f"[DocExtractor] OCR skipped after {max_attempts} attempt(s) for {label}: {exc}"
                )

    return ""


def _ocr_image_file(path: str) -> str:
    """
    Run pytesseract OCR directly on an image file path.

    Retries at most once on transient failure.  If the retry also fails,
    the image is silently skipped.
    """
    if not _check_tesseract():
        logger.warning(
            f"[DocExtractor] Tesseract OCR not available — skipping image {os.path.basename(path)}. "
            "Install Tesseract-OCR to enable image text extraction."
        )
        return ""

    max_attempts = 2
    last_exc: Optional[Exception] = None

    for attempt in range(1, max_attempts + 1):
        try:
            from PIL import Image
            import pytesseract
            _configure_tesseract_subprocess()
            img = Image.open(path)
            return pytesseract.image_to_string(img).strip()
        except Exception as exc:
            last_exc = exc
            if attempt < max_attempts:
                logger.warning(
                    f"[DocExtractor] OCR attempt {attempt}/{max_attempts} failed for "
                    f"{os.path.basename(path)}: {exc} — retrying"
                )
            else:
                logger.warning(
                    f"[DocExtractor] OCR skipped after {max_attempts} attempt(s) for "
                    f"{os.path.basename(path)}: {exc}"
                )

    return ""


def _is_duplicate(selectable_text: str, ocr_text: str) -> bool:
    """Check if the OCR extracted text is already contained/duplicated in the selectable text."""
    clean_sel = re.sub(r'[^a-z0-9]', '', selectable_text.lower())
    clean_ocr = re.sub(r'[^a-z0-9]', '', ocr_text.lower())
    if not clean_ocr:
        return True

    # If the clean OCR text is fully contained in selectable text
    if clean_ocr in clean_sel:
        return True

    # Calculate word-level overlap
    ocr_words = ocr_text.lower().split()
    if not ocr_words:
        return True
    sel_words = set(selectable_text.lower().split())
    overlap = sum(1 for w in ocr_words if w in sel_words)
    ratio = overlap / len(ocr_words)
    if ratio > 0.85:
        return True
    return False


def _extract_pdf(path: str) -> str:
    """Extract plain text and run OCR on embedded images in page order."""
    try:
        import fitz
        doc = fitz.open(path)
        parts = []

        for page_idx, page in enumerate(doc, 1):
            page_parts = []
            # Extract selectable text
            text = page.get_text()
            if text.strip():
                page_parts.append(text.strip())

            # Extract page images
            try:
                images = page.get_images(full=True)
                for img_info in images:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ocr_text = _ocr_image_bytes(
                        image_bytes,
                        label=f"PDF page {page_idx} img xref={xref}",
                    )
                    if ocr_text and not _is_duplicate(text, ocr_text):
                        page_parts.append(f"[Embedded Image OCR: {ocr_text}]")
            except Exception as img_err:
                logger.warning(
                    f"[DocExtractor] PDF image extraction failed on page {page_idx} of {path}: {img_err}"
                )

            if page_parts:
                parts.append("\n\n".join(page_parts))

        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"[DocExtractor] PDF PyMuPDF extraction failed for {path}: {e}. Falling back to pypdf.")
        # Fallback to pypdf
        try:
            from pypdf import PdfReader
            reader = PdfReader(path)
            parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    parts.append(text.strip())
            return "\n\n".join(parts)
        except Exception as pypdf_err:
            logger.warning(f"[DocExtractor] PDF fallback pypdf extraction also failed: {pypdf_err}")
            return ""


def _format_table_matrix_to_markdown(matrix: list[list[str]], title: str = "Table") -> str:
    r"""
    Format a 2D matrix of cell strings into a clean, structured Markdown table.
    Ensures:
      - All rows have the same number of columns (padded with empty strings).
      - Multi-line cell text (newlines) is converted to '<br>' to preserve single-line row structure.
      - Pipe characters '|' in cell content are escaped as '\|'.
      - Header separator row '| --- | --- | ... |' is generated.
      - Completely empty rows are omitted.
    """
    if not matrix:
        return ""

    cleaned_rows: list[list[str]] = []
    for row in matrix:
        cleaned_row = []
        for cell in row:
            c_str = str(cell).strip() if cell is not None else ""
            # Escape pipe characters
            c_str = c_str.replace("|", "\\|")
            # Replace internal newlines with <br>
            c_str = re.sub(r'\r?\n', '<br>', c_str)
            cleaned_row.append(c_str)
        # Keep row if at least one cell has non-whitespace text
        if any(c for c in cleaned_row):
            cleaned_rows.append(cleaned_row)

    if not cleaned_rows:
        return ""

    num_cols = max(len(r) for r in cleaned_rows)
    if num_cols == 0:
        return ""

    # Pad all rows to num_cols
    for row in cleaned_rows:
        while len(row) < num_cols:
            row.append("")

    lines = []
    header_row = cleaned_rows[0]
    lines.append(f"[{title}: {len(cleaned_rows)} rows x {num_cols} cols]")
    lines.append("| " + " | ".join(header_row) + " |")
    lines.append("| " + " | ".join(["---"] * num_cols) + " |")
    for row in cleaned_rows[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _process_docx_paragraph(para, doc, parts_acc: list[str]) -> None:
    """Extract paragraph text and OCR text from embedded drawings inside runs."""
    para_parts = []
    text = para.text.strip()
    if text:
        para_parts.append(text)

    # Check for drawings inside runs
    for run in para.runs:
        try:
            drawings = run._r.findall(
                './/{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'
            )
            for drawing in drawings:
                embeds = drawing.findall(
                    './/{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                )
                for embed in embeds:
                    rId = embed.get(
                        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                    )
                    if rId and hasattr(doc, "part") and rId in doc.part.related_parts:
                        image_part = doc.part.related_parts[rId]
                        image_bytes = image_part.blob
                        ocr_text = _ocr_image_bytes(
                            image_bytes, label=f"DOCX para embed rId={rId}"
                        )
                        if ocr_text and not _is_duplicate(text, ocr_text):
                            para_parts.append(f"[Embedded Image OCR: {ocr_text}]")
        except Exception:
            pass

    if para_parts:
        parts_acc.append("\n".join(para_parts))


def _process_docx_cell_text(cell, doc) -> str:
    """Extract cell text including nested tables and paragraph OCR drawings."""
    cell_parts = []
    # Loop over block-level items inside cell XML
    for elem in cell._tc:
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        if tag == "p":
            from docx.text.paragraph import Paragraph
            p = Paragraph(elem, cell)
            p_text = p.text.strip()
            if p_text:
                cell_parts.append(p_text)
            # OCR drawings in paragraph
            for run in p.runs:
                try:
                    drawings = run._r.findall(
                        './/{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'
                    )
                    for drawing in drawings:
                        embeds = drawing.findall(
                            './/{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                        )
                        for embed in embeds:
                            rId = embed.get(
                                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                            )
                            if rId and hasattr(doc, "part") and rId in doc.part.related_parts:
                                image_part = doc.part.related_parts[rId]
                                image_bytes = image_part.blob
                                ocr_text = _ocr_image_bytes(
                                    image_bytes, label=f"DOCX cell img {rId}"
                                )
                                if ocr_text and not _is_duplicate(p_text, ocr_text):
                                    cell_parts.append(f"[Embedded Image OCR: {ocr_text}]")
                except Exception:
                    pass
        elif tag == "tbl":
            from docx.table import Table
            nested_tbl = Table(elem, cell)
            nested_md = _process_docx_table(nested_tbl, doc, is_nested=True)
            if nested_md:
                cell_parts.append(f"\n[Nested Table]\n{nested_md}\n")

    return "\n".join(cell_parts).strip()


def _process_docx_table(table, doc, is_nested: bool = False) -> str:
    """Extract table as a Markdown table matrix with merged cell handling."""
    matrix = []
    for row in table.rows:
        row_cells = []
        seen_tc = set()
        for cell in row.cells:
            tc_id = id(cell._tc)
            # Horizontal merge in python-docx: adjacent cells in row.cells point to same _tc
            if tc_id in seen_tc:
                # Spanned cell (horizontally merged) — keep empty to maintain layout
                row_cells.append("")
                continue
            seen_tc.add(tc_id)

            c_text = _process_docx_cell_text(cell, doc)
            row_cells.append(c_text)

        if any(c for c in row_cells):
            matrix.append(row_cells)

    title = "Nested Table" if is_nested else "Table"
    return _format_table_matrix_to_markdown(matrix, title=title)


def _extract_docx(path: str) -> str:
    """Extract plain text and tables in document flow order, with OCR on embedded drawings."""
    try:
        from docx import Document
        from docx.text.paragraph import Paragraph
        from docx.table import Table

        doc = Document(path)
        parts = []

        # Iterate over child elements of document body in native document order
        for child in doc.element.body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "p":
                para = Paragraph(child, doc)
                _process_docx_paragraph(para, doc, parts)
            elif tag == "tbl":
                table = Table(child, doc)
                tbl_md = _process_docx_table(table, doc)
                if tbl_md:
                    parts.append(tbl_md)
            elif tag in ("sdt", "txbxContent"):
                # Structured Document Tags or Text Boxes containing paragraphs & tables
                for sub in child.iter():
                    sub_tag = sub.tag.split("}")[-1] if "}" in sub.tag else sub.tag
                    if sub_tag == "p":
                        p = Paragraph(sub, doc)
                        _process_docx_paragraph(p, doc, parts)
                    elif sub_tag == "tbl":
                        t = Table(sub, doc)
                        t_md = _process_docx_table(t, doc)
                        if t_md:
                            parts.append(t_md)

        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"[DocExtractor] DOCX extraction failed for {path}: {e}")
        return ""


def _process_pptx_shape(shape, slide_idx: int, slide_parts: list[str], extracted_set: set[str]) -> None:
    """Recursively extract content from a pptx shape."""
    # 1. Group shapes (recurse into sub-shapes)
    if hasattr(shape, "shapes"):
        try:
            for sub_shape in shape.shapes:
                _process_pptx_shape(sub_shape, slide_idx, slide_parts, extracted_set)
        except Exception as grp_err:
            logger.warning(f"[DocExtractor] Failed iterating sub-shapes of group on slide {slide_idx}: {grp_err}")
        return

    # 2. Table shape
    if hasattr(shape, "has_table") and shape.has_table:
        try:
            table = shape.table
            matrix = []
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    is_spanned = False
                    try:
                        if hasattr(cell, "is_spanned") and cell.is_spanned:
                            is_spanned = True
                    except Exception:
                        pass

                    cell_text = cell.text.strip() if hasattr(cell, "text") else ""
                    if is_spanned:
                        # Spanned cell in python-pptx — keep empty to maintain column alignment
                        row_cells.append("")
                    else:
                        row_cells.append(cell_text)

                matrix.append(row_cells)

            tbl_md = _format_table_matrix_to_markdown(matrix, title=f"Table (Slide {slide_idx})")
            if tbl_md:
                norm = re.sub(r'\s+', ' ', tbl_md.lower())
                if not any(norm in x or x in norm for x in extracted_set):
                    extracted_set.add(norm)
                    slide_parts.append(tbl_md)
        except Exception as tbl_err:
            logger.warning(f"[DocExtractor] Failed extracting table on slide {slide_idx}: {tbl_err}")
        return

    # 3. Chart shape
    if hasattr(shape, "has_chart") and shape.has_chart:
        try:
            chart = shape.chart
            chart_texts = []
            
            # Title
            if chart.has_title:
                try:
                    if chart.chart_title.has_text_frame and chart.chart_title.text_frame.text:
                        chart_texts.append(f"Title: {chart.chart_title.text_frame.text.strip()}")
                except Exception:
                    pass

            # Series & Categories
            try:
                for plot in chart.plots:
                    series_names = [s.name.strip() for s in plot.series if s.name]
                    if series_names:
                        chart_texts.append(f"Series: {', '.join(series_names)}")
                    if hasattr(plot, "categories"):
                        cat_labels = []
                        for cat in plot.categories:
                            if hasattr(cat, "label") and cat.label:
                                cat_labels.append(str(cat.label).strip())
                            elif cat:
                                cat_labels.append(str(cat).strip())
                        if cat_labels:
                            chart_texts.append(f"Categories: {', '.join(cat_labels)}")
            except Exception:
                pass

            # Category Axis Title
            try:
                if hasattr(chart, "category_axis") and chart.category_axis.has_title:
                    axis_title = chart.category_axis.axis_title.text_frame.text.strip()
                    if axis_title:
                        chart_texts.append(f"Category Axis Title: {axis_title}")
            except Exception:
                pass

            if chart_texts:
                chart_text = "[Chart]\n" + "\n".join(chart_texts)
                norm = re.sub(r'\s+', ' ', chart_text.lower())
                if not any(norm in x or x in norm for x in extracted_set):
                    extracted_set.add(norm)
                    slide_parts.append(chart_text)
        except Exception as chart_err:
            logger.warning(f"[DocExtractor] Failed extracting chart on slide {slide_idx}: {chart_err}")
        return

    # 4. Standard shape/Text frame
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        try:
            text = shape.text_frame.text.strip()
            if text:
                norm = re.sub(r'\s+', ' ', text.lower())
                if not any(norm in x or x in norm for x in extracted_set):
                    extracted_set.add(norm)
                    slide_parts.append(text)
        except Exception as txt_err:
            logger.warning(f"[DocExtractor] Failed extracting text frame on slide {slide_idx}: {txt_err}")
    elif hasattr(shape, "text") and shape.text.strip():
        try:
            text = shape.text.strip()
            norm = re.sub(r'\s+', ' ', text.lower())
            if not any(norm in x or x in norm for x in extracted_set):
                extracted_set.add(norm)
                slide_parts.append(text)
        except Exception as txt_err:
            logger.warning(f"[DocExtractor] Failed extracting text property on slide {slide_idx}: {txt_err}")

    # 5. Picture shape
    is_picture = False
    try:
        if getattr(shape, "shape_type", None) == 13 or hasattr(shape, "image"):
            is_picture = True
    except Exception:
        pass

    if is_picture:
        try:
            image_bytes = shape.image.blob
            ocr_text = _ocr_image_bytes(
                image_bytes, label=f"PPTX slide {slide_idx} shape {shape.shape_id}"
            )
            if ocr_text:
                norm = re.sub(r'\s+', ' ', ocr_text.lower())
                if not any(norm in x or x in norm for x in extracted_set):
                    extracted_set.add(norm)
                    slide_parts.append(f"[Embedded Image OCR: {ocr_text}]")
        except Exception as ocr_err:
            logger.warning(f"[DocExtractor] OCR failed on slide {slide_idx} image (shape_id={getattr(shape, 'shape_id', 'unknown')}): {ocr_err}")

    # 6. XML Text fallback (for SmartArt, drawing ML shapes, custom graphic objects, etc.)
    try:
        xml_texts = []
        if hasattr(shape, "element") and shape.element is not None:
            for node in shape.element.iter():
                if node.tag.endswith("}t") and node.text:
                    txt = node.text.strip()
                    if txt:
                        xml_texts.append(txt)
        if xml_texts:
            full_xml_text = " ".join(xml_texts)
            norm = re.sub(r'\s+', ' ', full_xml_text.lower())
            if not any(norm in x or x in norm for x in extracted_set):
                unique_xml_texts = []
                for t in xml_texts:
                    t_norm = re.sub(r'\s+', ' ', t.lower())
                    if not any(t_norm in x or x in t_norm for x in extracted_set):
                        unique_xml_texts.append(t)
                        extracted_set.add(t_norm)
                if unique_xml_texts:
                    slide_parts.append(" ".join(unique_xml_texts))
    except Exception as xml_err:
        logger.debug(f"[DocExtractor] XML text extraction failed on shape: {xml_err}")


def _extract_pptx(path: str) -> str:
    """Extract recursively all slide content (text, tables, charts, SmartArt, notes, pictures OCR) in order."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_parts = []
            extracted_set = set()

            # Process all slide shapes recursively
            for shape in slide.shapes:
                _process_pptx_shape(shape, slide_idx, slide_parts, extracted_set)

            # Process speaker notes
            try:
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        norm = re.sub(r'\s+', ' ', notes.lower())
                        if not any(norm in x or x in norm for x in extracted_set):
                            extracted_set.add(norm)
                            slide_parts.append(f"[Speaker Notes]\n{notes}")
            except Exception as notes_err:
                logger.debug(f"[DocExtractor] Speaker notes extraction failed on slide {slide_idx}: {notes_err}")

            # Combine slide parts (if empty, we still output the slide marker to preserve indices)
            slide_content = "\n\n".join(slide_parts)
            if slide_content:
                parts.append(f"[Slide {slide_idx}]\n{slide_content}")
            else:
                parts.append(f"[Slide {slide_idx}]")

        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"[DocExtractor] PPTX extraction failed for {path}: {e}")
        return ""


def _extract_text(path: str) -> str:
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            return f.read().strip()
    except Exception as e:
        logger.warning(f"[DocExtractor] Text file read failed for {path}: {e}")
        return ""


def _extract_image_ocr(path: str) -> str:
    """OCR a standalone image file (PNG, JPG, JPEG, WEBP)."""
    return _ocr_image_file(path)


def _format_dataframe(df: str, filename: str, sheet_name: str) -> str:
    import pandas as pd
    # df is pd.DataFrame
    # Drop completely empty rows and columns
    df = df.dropna(how="all").dropna(how="all", axis=1)
    if df.empty:
        return ""

    df = df.fillna("")
    cols = [str(c).strip() for c in df.columns]

    block_size = 20
    formatted_blocks = []

    for start_idx in range(0, len(df), block_size):
        end_idx = min(start_idx + block_size, len(df))
        df_block = df.iloc[start_idx:end_idx]

        row_numbers = [idx + 2 for idx in df_block.index]
        first_row = row_numbers[0]
        last_row = row_numbers[-1]

        block_lines = [
            f"Source: {filename}",
            f"Sheet: {sheet_name}",
            f"Rows: {first_row}-{last_row}",
            "---------------------"
        ]

        header_line = " | ".join(cols)
        block_lines.append(f"Headers: {header_line}")

        for idx, row in df_block.iterrows():
            row_num = idx + 2
            row_vals = [str(val).strip() for val in row.values]
            row_str = " | ".join(row_vals)
            block_lines.append(f"[Row {row_num}] {row_str}")

        formatted_blocks.append("\n".join(block_lines))

    return "\n\n".join(formatted_blocks)


def _extract_spreadsheet(path: str, filename: str) -> str:
    """Extract worksheets from excel sheets/csv tables, partitioning rows with tracking metadata."""
    import pandas as pd
    ext = os.path.splitext(filename.lower())[1]
    parts = []

    try:
        if ext == ".csv":
            sheet_name = os.path.splitext(filename)[0]
            df = pd.read_csv(path)
            formatted = _format_dataframe(df, filename, sheet_name)
            if formatted:
                parts.append(formatted)
        else:
            excel_file = pd.ExcelFile(path)
            for sheet_name in excel_file.sheet_names:
                df = excel_file.parse(sheet_name)
                formatted = _format_dataframe(df, filename, sheet_name)
                if formatted:
                    parts.append(formatted)

        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"[DocExtractor] Spreadsheet extraction failed for {filename}: {e}")
        return ""


def _get_soffice_cmd() -> str | None:
    """Return path or binary name for LibreOffice soffice command if available."""
    import shutil
    # 1. Check system PATH
    soffice_path = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice_path:
        return soffice_path

    # 2. Check standard Windows installations
    if sys.platform == "win32":
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for cand in candidates:
            if os.path.isfile(cand):
                return cand
    return None


def _convert_legacy_doc_or_ppt(file_path: str, target_ext: str) -> str | None:
    """
    Convert legacy binary formats (.doc -> .docx, .ppt -> .pptx) using LibreOffice headless mode.
    Returns the path to the temporary converted file, or None if conversion failed/unavailable.
    """
    soffice_cmd = _get_soffice_cmd()
    if not soffice_cmd:
        logger.warning(
            f"[DocExtractor] LibreOffice (soffice) not found on system. "
            f"Cannot convert legacy '{os.path.basename(file_path)}' to {target_ext}. "
            "Please install LibreOffice or convert file to .docx/.pptx."
        )
        return None

    import tempfile
    temp_dir = tempfile.mkdtemp(prefix="voicesum_doc_conv_")
    try:
        cmd = [
            soffice_cmd,
            "--headless",
            "--convert-to",
            target_ext.lstrip("."),
            "--outdir",
            temp_dir,
            file_path,
        ]
        kw = {}
        if sys.platform == "win32":
            kw["creationflags"] = _CREATE_NO_WINDOW
        proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60, **kw)
        if proc.returncode == 0:
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            converted_path = os.path.join(temp_dir, base_name + target_ext)
            if os.path.isfile(converted_path):
                return converted_path
    except Exception as exc:
        logger.warning(f"[DocExtractor] Legacy document conversion failed for {file_path}: {exc}")

    return None


SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md",
    ".png", ".jpg", ".jpeg", ".webp",
    ".xlsx", ".xls", ".csv",
}

UNSUPPORTED_EXTENSIONS: set[str] = set()


def extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract plain text from a document, image, or spreadsheet file."""
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        text = _extract_pdf(file_path)
    elif ext == ".docx":
        text = _extract_docx(file_path)
    elif ext == ".pptx":
        text = _extract_pptx(file_path)
    elif ext == ".doc":
        converted = _convert_legacy_doc_or_ppt(file_path, ".docx")
        if converted and os.path.isfile(converted):
            try:
                text = _extract_docx(converted)
            finally:
                try:
                    os.remove(converted)
                    os.rmdir(os.path.dirname(converted))
                except Exception:
                    pass
        else:
            text = ""
    elif ext == ".ppt":
        converted = _convert_legacy_doc_or_ppt(file_path, ".pptx")
        if converted and os.path.isfile(converted):
            try:
                text = _extract_pptx(converted)
            finally:
                try:
                    os.remove(converted)
                    os.rmdir(os.path.dirname(converted))
                except Exception:
                    pass
        else:
            text = ""
    elif ext in (".txt", ".md"):
        text = _extract_text(file_path)
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        text = _extract_image_ocr(file_path)
    elif ext in (".xlsx", ".xls", ".csv"):
        text = _extract_spreadsheet(file_path, filename)
    else:
        logger.warning(f"[DocExtractor] Unknown extension '{ext}' for '{filename}' — skipping.")
        return ""

    logger.info(f"[DocExtractor] Extracted {len(text)} chars from '{filename}' (format={ext})")
    return text


def compute_file_hash(file_path: str) -> str:
    """Compute SHA-256 hash of a file for change detection."""
    import hashlib
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()
